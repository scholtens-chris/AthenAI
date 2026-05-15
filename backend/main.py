import os
import logging
from logging.handlers import RotatingFileHandler
from fastapi import FastAPI, UploadFile, File, Form, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from typing import List, Optional
from contextlib import asynccontextmanager
from collections import Counter
from pathlib import Path

import uuid
import zipfile
import io
import json
import xml.etree.ElementTree as ET
import re
import time

from qwen_utils import QwenLLM


def configure_logging() -> logging.Logger:
    log_level = os.getenv("ATHENAI_LOG_LEVEL", "INFO").upper()
    log_format = "%(asctime)s %(levelname)s [%(name)s] %(message)s"
    root_logger = logging.getLogger()
    if not root_logger.handlers:
        logging.basicConfig(level=log_level, format=log_format)
    else:
        root_logger.setLevel(log_level)

    log_file = os.getenv("ATHENAI_LOG_FILE")
    if log_file and not any(
        isinstance(handler, RotatingFileHandler) and getattr(handler, "baseFilename", None) == str(Path(log_file).resolve())
        for handler in root_logger.handlers
    ):
        file_handler = RotatingFileHandler(log_file, maxBytes=5_000_000, backupCount=3, encoding="utf-8")
        file_handler.setLevel(log_level)
        file_handler.setFormatter(logging.Formatter(log_format))
        root_logger.addHandler(file_handler)

    return logging.getLogger("athenai.api")


logger = configure_logging()

# In-memory session store (replace with Redis for production)
sessions = {}

CHUNK_WORDS = int(os.getenv("ATHENAI_RAG_CHUNK_WORDS", "220"))
CHUNK_OVERLAP = int(os.getenv("ATHENAI_RAG_CHUNK_OVERLAP", "45"))
DEFAULT_RAG_TOP_K = int(os.getenv("ATHENAI_RAG_TOP_K", "4"))
MAX_CONTEXT_CHARS = int(os.getenv("ATHENAI_RAG_MAX_CONTEXT_CHARS", "8000"))
OVERVIEW_CHUNK_LIMIT = int(os.getenv("ATHENAI_RAG_OVERVIEW_CHUNKS", "8"))
FULL_ARTIFACT_CONTEXT_CHARS = int(os.getenv("ATHENAI_RAG_FULL_ARTIFACT_CONTEXT_CHARS", "24000"))
SESSION_TTL_SECONDS = int(float(os.getenv("ATHENAI_SESSION_TTL_HOURS", "12")) * 60 * 60)
SESSION_STORE_DIR = Path(os.getenv("ATHENAI_SESSION_STORE_DIR", Path(__file__).parent / ".athenai_sessions"))
SESSION_PERSISTENCE_ENABLED = os.getenv("ATHENAI_SESSION_PERSISTENCE", "1") != "0"

STOPWORDS = {
    "a", "an", "and", "are", "as", "at", "be", "but", "by", "for", "from", "how",
    "i", "in", "into", "is", "it", "of", "on", "or", "that", "the", "this", "to",
    "was", "what", "when", "where", "which", "who", "why", "with", "you", "your",
}

STUDY_TASK_PATTERNS = {
    "quiz": (
        r"\bquiz\b",
        r"\btest me\b",
        r"\bpractice question",
        r"\bflashcards?\b",
    ),
    "summary": (
        r"\bsummar(?:y|ize|ise|ise)\b",
        r"\brecap\b",
        r"\bchapter\s+\d+\b",
    ),
    "find": (
        r"\bfind\b",
        r"\bwhere\b",
        r"\bwhen\b",
        r"\btimestamp\b",
        r"\bexplained?\b",
        r"\bmentioned?\b",
    ),
    "takeaways": (
        r"\bkey takeaways?\b",
        r"\bmain points?\b",
        r"\bbig ideas?\b",
        r"\bimportant concepts?\b",
    ),
    "study_guide": (
        r"\bstudy\s+guides?\b",
        r"\bstudyguide\b",
        r"\breview\s+guides?\b",
        r"\blearning\s+objectives?\b",
    ),
    "explain": (
        r"\bexplain\b",
        r"\bteach me\b",
        r"\bwalk me through\b",
        r"\bhelp me understand\b",
        r"\blecture\b",
    ),
}

BROAD_STUDY_PROMPTS = {
    "quiz",
    "quiz me",
    "test me",
    "explain this",
    "explain this lecture",
    "summarize",
    "summarize this",
    "summarize this lecture",
    "what are key takeaways",
    "key takeaways",
    "main points",
    "study guide",
    "make a study guide",
    "create a study guide",
    "review guide",
}

# Initialize Qwen LLM (update model_path/device as needed)
qwen_llm = QwenLLM()


def session_now() -> int:
    return int(time.time())


def session_expires_at() -> int:
    return session_now() + max(60, SESSION_TTL_SECONDS)


def session_file_path(session_id: str) -> Path:
    safe_id = re.sub(r"[^a-zA-Z0-9_.-]", "_", session_id)
    return SESSION_STORE_DIR / f"{safe_id}.json"


def normalize_session(session: Optional[dict] = None) -> dict:
    session = session or {}
    session.setdefault("content", [])
    session.setdefault("chunks", [])
    session.setdefault("chat_history", [])
    session["updated_at"] = int(session.get("updated_at") or session_now())
    session["expires_at"] = int(session.get("expires_at") or session_expires_at())
    for chunk in session.get("chunks", []):
        if not isinstance(chunk.get("tokens"), Counter):
            chunk["tokens"] = Counter(chunk.get("tokens", {}))
    return session


def is_session_expired(session: dict) -> bool:
    return int(session.get("expires_at") or 0) <= session_now()


def touch_session(session_id: str) -> None:
    session = normalize_session(sessions.get(session_id))
    session["updated_at"] = session_now()
    session["expires_at"] = session_expires_at()
    sessions[session_id] = session


def persist_session(session_id: str) -> None:
    if not SESSION_PERSISTENCE_ENABLED:
        logger.debug("session_persist_skipped persistence_disabled session_id=%s", session_id)
        return
    session = sessions.get(session_id)
    if not session:
        logger.debug("session_persist_skipped empty_session session_id=%s", session_id)
        return
    SESSION_STORE_DIR.mkdir(parents=True, exist_ok=True)
    path = session_file_path(session_id)
    temp_path = path.with_suffix(".tmp")
    payload = normalize_session(session).copy()
    try:
        temp_path.write_text(json.dumps(payload), encoding="utf-8")
        temp_path.replace(path)
        logger.debug(
            "session_persisted session_id=%s path=%s chunks=%s messages=%s expires_at=%s",
            session_id,
            path,
            len(payload.get("chunks", [])),
            len(payload.get("chat_history", [])),
            payload.get("expires_at"),
        )
    except OSError:
        logger.exception("session_persist_failed session_id=%s path=%s", session_id, path)


def load_session_from_disk(session_id: str) -> Optional[dict]:
    if not SESSION_PERSISTENCE_ENABLED:
        return None
    path = session_file_path(session_id)
    if not path.exists():
        return None
    try:
        session = normalize_session(json.loads(path.read_text(encoding="utf-8")))
    except (OSError, json.JSONDecodeError):
        logger.exception("session_load_failed session_id=%s path=%s", session_id, path)
        return None
    if is_session_expired(session):
        sessions.pop(session_id, None)
        try:
            path.unlink()
        except OSError:
            logger.exception("session_expired_delete_failed session_id=%s path=%s", session_id, path)
        logger.info("session_expired session_id=%s path=%s", session_id, path)
        return None
    sessions[session_id] = session
    logger.info(
        "session_loaded session_id=%s path=%s chunks=%s messages=%s",
        session_id,
        path,
        len(session.get("chunks", [])),
        len(session.get("chat_history", [])),
    )
    return session


def cleanup_expired_sessions() -> None:
    removed_memory = 0
    for session_id, session in list(sessions.items()):
        if is_session_expired(session):
            sessions.pop(session_id, None)
            removed_memory += 1
    if not SESSION_PERSISTENCE_ENABLED or not SESSION_STORE_DIR.exists():
        if removed_memory:
            logger.info("session_cleanup_complete removed_memory=%s removed_disk=0", removed_memory)
        return
    removed_disk = 0
    for path in SESSION_STORE_DIR.glob("*.json"):
        try:
            session = json.loads(path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            logger.exception("session_cleanup_read_failed path=%s", path)
            continue
        if is_session_expired(session):
            try:
                path.unlink()
                removed_disk += 1
            except OSError:
                logger.exception("session_cleanup_delete_failed path=%s", path)
    if removed_memory or removed_disk:
        logger.info("session_cleanup_complete removed_memory=%s removed_disk=%s", removed_memory, removed_disk)


@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info(
        "startup_begin mock_llm=%s model=%s persistence=%s session_store=%s ttl_seconds=%s",
        qwen_llm.use_mock_model,
        qwen_llm.model_path,
        SESSION_PERSISTENCE_ENABLED,
        SESSION_STORE_DIR,
        SESSION_TTL_SECONDS,
    )
    cleanup_expired_sessions()
    preload_model = os.getenv("ATHENAI_PRELOAD_MODEL", "0") == "1"
    if preload_model and not qwen_llm.use_mock_model:
        logger.info("model_preload_begin model=%s device=%s", qwen_llm.model_path, qwen_llm.device)
        qwen_llm._load_model()
        logger.info("model_preload_complete model=%s device=%s", qwen_llm.model_path, qwen_llm.device)
    yield
    logger.info("shutdown_complete")

app = FastAPI(lifespan=lifespan)


@app.middleware("http")
async def log_requests(request: Request, call_next):
    request_id = request.headers.get("x-request-id") or str(uuid.uuid4())
    start = time.perf_counter()
    logger.info("request_begin request_id=%s method=%s path=%s", request_id, request.method, request.url.path)
    try:
        response = await call_next(request)
    except Exception:
        elapsed_ms = int((time.perf_counter() - start) * 1000)
        logger.exception(
            "request_failed request_id=%s method=%s path=%s elapsed_ms=%s",
            request_id,
            request.method,
            request.url.path,
            elapsed_ms,
        )
        raise
    elapsed_ms = int((time.perf_counter() - start) * 1000)
    response.headers["x-request-id"] = request_id
    logger.info(
        "request_complete request_id=%s method=%s path=%s status=%s elapsed_ms=%s",
        request_id,
        request.method,
        request.url.path,
        response.status_code,
        elapsed_ms,
    )
    return response

# Enable CORS for frontend integration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def parse_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="ignore")

def ensure_session(session_id: str):
    if session_id not in sessions:
        load_session_from_disk(session_id)
    if session_id not in sessions:
        sessions[session_id] = normalize_session()
    sessions[session_id] = normalize_session(sessions[session_id])
    if is_session_expired(sessions[session_id]):
        sessions[session_id] = normalize_session()
    return sessions[session_id]

def tokenize(text: str) -> list[str]:
    return [
        token
        for token in re.findall(r"[a-zA-Z0-9']+", text.lower())
        if len(token) > 1 and token not in STOPWORDS
    ]

def chunk_text(text: str, chunk_words: int = CHUNK_WORDS, overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = re.findall(r"\S+", text)
    if not words:
        return []
    if len(words) <= chunk_words:
        return [" ".join(words)]

    chunks = []
    step = max(1, chunk_words - overlap)
    for start in range(0, len(words), step):
        chunk = " ".join(words[start:start + chunk_words]).strip()
        if chunk:
            chunks.append(chunk)
        if start + chunk_words >= len(words):
            break
    return chunks

def add_document_to_session(session_id: str, filename: str, text: str) -> int:
    session = ensure_session(session_id)
    document_id = str(uuid.uuid4())
    cleaned_text = re.sub(r"\s+", " ", text).strip()
    chunks = chunk_text(cleaned_text)
    if not chunks:
        logger.warning("document_skipped_no_text session_id=%s filename=%s", session_id, filename)
        return 0

    session["content"].append({"id": document_id, "filename": filename, "text": cleaned_text})

    for index, chunk in enumerate(chunks):
        tokens = tokenize(chunk)
        session["chunks"].append({
            "id": f"{document_id}:{index}",
            "document_id": document_id,
            "filename": filename,
            "chunk_index": index,
            "text": chunk,
            "tokens": Counter(tokens),
            "token_count": max(1, len(tokens)),
        })
    touch_session(session_id)
    persist_session(session_id)
    logger.info(
        "document_indexed session_id=%s filename=%s chars=%s chunks_added=%s session_chunks=%s",
        session_id,
        filename,
        len(cleaned_text),
        len(chunks),
        len(session["chunks"]),
    )
    return len(chunks)

def retrieve_context(session_id: Optional[str], query: str, top_k: int = DEFAULT_RAG_TOP_K) -> list[dict]:
    if not session_id:
        logger.debug("retrieve_context_skipped no_session query_chars=%s", len(query or ""))
        return []
    session = ensure_session(session_id)
    if not session.get("chunks"):
        logger.info("retrieve_context_empty_session session_id=%s query_chars=%s", session_id, len(query or ""))
        return []

    query_tokens = tokenize(query)
    if not query_tokens:
        logger.info("retrieve_context_empty_query session_id=%s query_chars=%s", session_id, len(query or ""))
        return []

    query_counts = Counter(query_tokens)
    scored = []
    for chunk in session.get("chunks", []):
        overlap_terms = set(query_counts) & set(chunk["tokens"])
        if not overlap_terms:
            continue
        term_score = sum(min(query_counts[token], chunk["tokens"][token]) for token in overlap_terms)
        coverage_score = len(overlap_terms) / max(1, len(query_counts))
        density_score = term_score / chunk["token_count"]
        score = term_score + (coverage_score * 2.0) + density_score
        scored.append((score, chunk))

    scored.sort(key=lambda item: item[0], reverse=True)

    results = []
    used_chars = 0
    for score, chunk in scored[:max(1, top_k)]:
        text = chunk["text"]
        if used_chars + len(text) > MAX_CONTEXT_CHARS:
            text = text[:max(0, MAX_CONTEXT_CHARS - used_chars)].rstrip()
        if not text:
            break
        used_chars += len(text)
        results.append({
            "id": chunk["id"],
            "filename": chunk["filename"],
            "chunk_index": chunk["chunk_index"],
            "score": round(score, 4),
            "text": text,
        })
        if used_chars >= MAX_CONTEXT_CHARS:
            break
    logger.info(
        "retrieve_context_complete session_id=%s query_tokens=%s scored_chunks=%s returned=%s used_chars=%s top_k=%s",
        session_id,
        len(query_tokens),
        len(scored),
        len(results),
        used_chars,
        top_k,
    )
    return results

def infer_study_task(prompt: str) -> str:
    normalized = re.sub(r"\s+", " ", prompt.lower()).strip()
    for task, patterns in STUDY_TASK_PATTERNS.items():
        if any(re.search(pattern, normalized) for pattern in patterns):
            return task
    return "answer"

def is_broad_study_prompt(prompt: str, task: str) -> bool:
    normalized = re.sub(r"\s+", " ", prompt.lower()).strip(" ?.!")
    content_tokens = tokenize(normalized)
    if normalized in BROAD_STUDY_PROMPTS:
        return True
    full_material_terms = {
        "all",
        "everything",
        "entire",
        "full",
        "uploaded",
        "material",
        "materials",
        "notes",
        "lecture",
        "lectures",
        "document",
        "documents",
        "file",
        "files",
    }
    if set(content_tokens) & full_material_terms:
        return True
    if task in {"quiz", "takeaways", "study_guide"} and len(content_tokens) <= 3:
        return True
    return False

def retrieve_session_overview(session_id: Optional[str], limit: int = OVERVIEW_CHUNK_LIMIT) -> list[dict]:
    if not session_id:
        return []

    chunks = ensure_session(session_id).get("chunks", [])
    if not chunks:
        return []

    selected = []
    seen_documents = set()
    for chunk in chunks:
        document_id = chunk["document_id"]
        if document_id in seen_documents:
            continue
        seen_documents.add(document_id)
        selected.append(chunk)
        if len(selected) >= limit:
            break

    if len(selected) < limit:
        selected_ids = {chunk["id"] for chunk in selected}
        stride = max(1, len(chunks) // max(1, limit))
        for index in range(0, len(chunks), stride):
            chunk = chunks[index]
            if chunk["id"] not in selected_ids:
                selected.append(chunk)
                selected_ids.add(chunk["id"])
            if len(selected) >= limit:
                break

    results = []
    used_chars = 0
    for chunk in selected[:limit]:
        text = chunk["text"]
        if used_chars + len(text) > MAX_CONTEXT_CHARS:
            text = text[:max(0, MAX_CONTEXT_CHARS - used_chars)].rstrip()
        if not text:
            break
        used_chars += len(text)
        results.append({
            "id": chunk["id"],
            "filename": chunk["filename"],
            "chunk_index": chunk["chunk_index"],
            "score": 0,
            "text": text,
        })
        if used_chars >= MAX_CONTEXT_CHARS:
            break
    return results

def material_context_entry(chunk: dict, score: float, text: str) -> dict:
    return {
        "id": chunk["id"],
        "filename": chunk["filename"],
        "chunk_index": chunk["chunk_index"],
        "score": score,
        "text": text,
    }

def retrieve_full_session_context(session_id: Optional[str], max_chars: int = FULL_ARTIFACT_CONTEXT_CHARS) -> list[dict]:
    if not session_id:
        logger.debug("retrieve_full_context_skipped no_session")
        return []

    chunks = ensure_session(session_id).get("chunks", [])
    if not chunks:
        logger.info("retrieve_full_context_empty session_id=%s", session_id)
        return []

    results = []
    used_chars = 0
    for chunk in chunks:
        text = chunk["text"]
        if max_chars > 0 and used_chars + len(text) > max_chars:
            remaining_chars = max(0, max_chars - used_chars)
            text = text[:remaining_chars].rstrip()
        if not text:
            break
        used_chars += len(text)
        results.append(material_context_entry(chunk, 0, text))
        if max_chars > 0 and used_chars >= max_chars:
            break
    logger.info(
        "retrieve_full_context_complete session_id=%s source_chunks=%s returned=%s used_chars=%s max_chars=%s",
        session_id,
        len(chunks),
        len(results),
        used_chars,
        max_chars,
    )
    return results

def retrieve_study_context(session_id: Optional[str], prompt: str, top_k: int) -> tuple[list[dict], str]:
    task = infer_study_task(prompt)
    if is_broad_study_prompt(prompt, task):
        context = retrieve_full_session_context(session_id)
        logger.info(
            "retrieve_study_context_full session_id=%s task=%s prompt_chars=%s context_chunks=%s max_chars=%s",
            session_id,
            task,
            len(prompt or ""),
            len(context),
            FULL_ARTIFACT_CONTEXT_CHARS,
        )
        return context, task

    context = retrieve_context(session_id, prompt, top_k)
    logger.info(
        "retrieve_study_context_retrieved session_id=%s task=%s prompt_chars=%s context_chunks=%s top_k=%s",
        session_id,
        task,
        len(prompt or ""),
        len(context),
        top_k,
    )
    return context, task

def public_source(source: dict) -> dict:
    return {
        "id": source["id"],
        "filename": source["filename"],
        "chunk_index": source["chunk_index"],
        "score": source["score"],
        "preview": source["text"][:360],
    }


def public_chat_message(message: dict) -> dict:
    return {
        key: value
        for key, value in message.items()
        if key in {"sender", "text", "sources", "tokenUsage", "examples"}
    }


def append_chat_message(session_id: Optional[str], message: dict) -> None:
    if not session_id:
        return
    session = ensure_session(session_id)
    session.setdefault("chat_history", []).append(message)
    touch_session(session_id)
    persist_session(session_id)


@app.get("/llm/status")
async def llm_status():
    """Report the currently configured LLM mode and model."""
    return {
        "model": qwen_llm.model_path,
        "mock": qwen_llm.use_mock_model,
        "loaded": qwen_llm.model is not None,
        "device": str(qwen_llm.device),
        "max_new_tokens": qwen_llm.max_new_tokens,
        "long_task_max_new_tokens": qwen_llm.long_task_max_new_tokens,
        "temperature": qwen_llm.temperature,
        "top_p": qwen_llm.top_p,
        "top_k": qwen_llm.top_k,
    }


@app.get("/session/{session_id}")
async def get_session(session_id: str):
    """Restore short-lived chat/session state after a UI or API restart."""
    session = ensure_session(session_id)
    return {
        "session_id": session_id,
        "expires_at": session["expires_at"],
        "chat_history": [public_chat_message(message) for message in session.get("chat_history", [])],
        "documents": [
            {"id": item["id"], "filename": item["filename"]}
            for item in session.get("content", [])
        ],
        "chunk_count": len(session.get("chunks", [])),
    }


def parse_json(content: bytes) -> str:
    try:
        data = json.loads(content.decode("utf-8", errors="ignore"))
        if isinstance(data, dict):
            return json.dumps(data, indent=2)
        elif isinstance(data, list):
            return "\n".join([json.dumps(item) for item in data])
        else:
            return str(data)
    except Exception:
        logger.exception("parse_json_failed bytes=%s", len(content or b""))
        return "[Invalid JSON]"

def parse_dfxp(content: bytes) -> str:
    try:
        root = ET.fromstring(content.decode("utf-8", errors="ignore"))
        texts = [elem.text for elem in root.iter() if elem.text]
        return " ".join(texts)
    except Exception:
        logger.exception("parse_dfxp_failed bytes=%s", len(content or b""))
        return "[Invalid DFXP XML]"

def parse_srt(content: bytes) -> str:
    text = content.decode("utf-8", errors="ignore")
    # Remove SRT numbering and timestamps
    text = re.sub(r"\d+\n\d{2}:\d{2}:\d{2},\d{3} --> \d{2}:\d{2}:\d{2},\d{3}", "", text)
    text = re.sub(r"\d{2}:\d{2}:\d{2}\.\d{3} --> \d{2}:\d{2}:\d{2}\.\d{3}", "", text)
    return re.sub(r"\n+", " ", text).strip()

def parse_vtt(content: bytes) -> str:
    text = content.decode("utf-8", errors="ignore")
    # Remove WEBVTT header and timestamps
    text = re.sub(r"WEBVTT.*?\n", "", text, flags=re.DOTALL)
    text = re.sub(r"\d{2}:\d{2}:\d{2}\.\d{3} --> \d{2}:\d{2}:\d{2}\.\d{3}", "", text)
    return re.sub(r"\n+", " ", text).strip()

def parse_pdf(content: bytes) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        logger.exception("parse_pdf_failed bytes=%s", len(content or b""))
        return ""

def parse_docx(content: bytes) -> str:
    try:
        from docx import Document

        document = Document(io.BytesIO(content))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    except Exception:
        logger.exception("parse_docx_failed bytes=%s", len(content or b""))
        return ""

def parse_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(content))
        texts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        logger.exception("parse_pptx_failed bytes=%s", len(content or b""))
        return ""

def extract_files_from_zip(zip_bytes: bytes):
    files = []
    with zipfile.ZipFile(io.BytesIO(zip_bytes)) as z:
        for name in z.namelist():
            if name.endswith("/"):
                continue
            files.append({"filename": name, "content": z.read(name)})
    logger.info("zip_extracted files=%s bytes=%s", len(files), len(zip_bytes or b""))
    return files

def parse_file(filename: str, content: bytes) -> str:
    ext = filename.lower().split(".")[-1]
    logger.debug("parse_file_begin filename=%s ext=%s bytes=%s", filename, ext, len(content or b""))
    if ext in {"txt", "md", "csv"}:
        return parse_txt(content)
    elif ext == "json":
        return parse_json(content)
    elif ext == "dfxp":
        return parse_dfxp(content)
    elif ext == "srt":
        return parse_srt(content)
    elif ext == "vtt":
        return parse_vtt(content)
    elif ext == "pdf":
        return parse_pdf(content)
    elif ext == "docx":
        return parse_docx(content)
    elif ext == "pptx":
        return parse_pptx(content)
    else:
        return content.decode("utf-8", errors="ignore")

@app.post("/upload/")
async def upload_content(session_id: Optional[str] = Form(None), files: List[UploadFile] = File(...)):
    """Bulk upload transcripts, OCR, and descriptions. Every upload is indexed for RAG."""
    if not session_id:
        session_id = str(uuid.uuid4())
    ensure_session(session_id)
    logger.info("upload_begin session_id=%s file_count=%s", session_id, len(files))
    indexed_files = []
    skipped_files = []
    added_chunk_count = 0
    for file in files:
        fname = file.filename
        content = await file.read()
        logger.info("upload_file_read session_id=%s filename=%s bytes=%s", session_id, fname, len(content or b""))
        if fname.lower().endswith(".zip"):
            try:
                extracted = extract_files_from_zip(content)
            except zipfile.BadZipFile:
                logger.exception("upload_zip_invalid session_id=%s filename=%s", session_id, fname)
                skipped_files.append(fname)
                continue
            for f in extracted:
                text = parse_file(f["filename"], f["content"])
                chunks_added = add_document_to_session(session_id, f["filename"], text)
                if chunks_added:
                    indexed_files.append(f["filename"])
                    added_chunk_count += chunks_added
                else:
                    skipped_files.append(f["filename"])
        else:
            text = parse_file(fname, content)
            chunks_added = add_document_to_session(session_id, fname, text)
            if chunks_added:
                indexed_files.append(fname)
                added_chunk_count += chunks_added
            else:
                skipped_files.append(fname)
    touch_session(session_id)
    persist_session(session_id)
    logger.info(
        "upload_complete session_id=%s uploaded_files=%s indexed=%s skipped=%s added_chunks=%s total_chunks=%s",
        session_id,
        len(files),
        len(indexed_files),
        len(skipped_files),
        added_chunk_count,
        len(sessions[session_id]["chunks"]),
    )
    return {
        "session_id": session_id,
        "files": [f.filename for f in files],
        "indexed_files": indexed_files,
        "skipped_files": skipped_files,
        "added_chunk_count": added_chunk_count,
        "chunk_count": len(sessions[session_id]["chunks"]),
    }


@app.post("/chat/")
async def chat(request: Request):
    """RAG-first chat endpoint for LLM interaction (Qwen)."""
    data = await request.json()
    prompt = data.get("prompt", "")
    session_id = data.get("session_id")
    use_internet = data.get("use_internet", False)
    image = data.get("image")  # Expecting base64 or bytes (optional)
    if not session_id:
        session_id = str(uuid.uuid4())
    ensure_session(session_id)
    logger.info(
        "chat_begin session_id=%s prompt_chars=%s use_internet=%s has_image=%s",
        session_id,
        len(prompt or ""),
        use_internet,
        image is not None,
    )

    top_k = int(data.get("top_k", DEFAULT_RAG_TOP_K))
    context, study_task = retrieve_study_context(session_id, prompt, top_k)
    logger.info(
        "chat_context session_id=%s study_task=%s context_chunks=%s top_k=%s",
        session_id,
        study_task,
        len(context),
        top_k,
    )
    append_chat_message(session_id, {"sender": "user", "text": prompt})
    if session_id and ensure_session(session_id).get("chunks") and not context:
        response = "I could not find relevant uploaded material for that question. Try asking with terms from your notes or upload more source files."
        append_chat_message(session_id, {"sender": "ai", "text": response, "sources": []})
        logger.warning("chat_no_relevant_context session_id=%s study_task=%s", session_id, study_task)
        return {
            "session_id": session_id,
            "response": response,
            "usage": qwen_llm.estimate_usage(prompt, response, context=[], study_task=study_task),
            "used_internet": False,
            "rag_first": True,
            "sources": [],
        }
    if not context:
        response = "Upload study files first so I can answer from your material."
        append_chat_message(session_id, {"sender": "ai", "text": response, "sources": []})
        logger.warning("chat_no_uploaded_context session_id=%s study_task=%s", session_id, study_task)
        return {
            "session_id": session_id,
            "response": response,
            "usage": qwen_llm.estimate_usage(prompt, response, context=[], study_task=study_task),
            "used_internet": False,
            "rag_first": True,
            "sources": [],
        }

    # TODO: handle image input (decode if base64)
    try:
        chat_result = qwen_llm.chat_with_usage(
            prompt,
            image=image,
            context=context,
            use_internet=use_internet,
            study_task=study_task,
        )
    except Exception:
        logger.exception(
            "chat_llm_failed session_id=%s study_task=%s context_chunks=%s",
            session_id,
            study_task,
            len(context),
        )
        raise
    sources = [public_source(source) for source in context]
    append_chat_message(
        session_id,
        {
            "sender": "ai",
            "text": chat_result["response"],
            "sources": sources,
            "tokenUsage": {
                "label": "Answer",
                "tokens": chat_result["usage"].get("completion_tokens"),
                "total": chat_result["usage"].get("total_tokens"),
                "maxNewTokens": chat_result["usage"].get("max_new_tokens"),
                "retriedForQuality": chat_result["usage"].get("retried_for_quality"),
                "hitTokenLimit": chat_result["usage"].get("hit_token_limit"),
                "estimated": chat_result["usage"].get("estimated"),
            },
        },
    )
    usage = chat_result.get("usage", {})
    logger.info(
        "chat_complete session_id=%s study_task=%s response_chars=%s prompt_tokens=%s completion_tokens=%s total_tokens=%s hit_token_limit=%s retried_for_quality=%s",
        session_id,
        study_task,
        len(chat_result.get("response", "")),
        usage.get("prompt_tokens"),
        usage.get("completion_tokens"),
        usage.get("total_tokens"),
        usage.get("hit_token_limit"),
        usage.get("retried_for_quality"),
    )
    return {
        "session_id": session_id,
        "response": chat_result["response"],
        "usage": chat_result["usage"],
        "used_internet": False,
        "rag_first": True,
        "study_task": study_task,
        "sources": sources,
    }

@app.post("/search/")
async def search(request: Request):
    """Search within uploaded content."""
    data = await request.json()
    session_id = data.get("session_id")
    query = data.get("query") or data.get("prompt") or ""
    top_k = int(data.get("top_k", DEFAULT_RAG_TOP_K))
    results = [public_source(source) for source in retrieve_context(session_id, query, top_k)]
    logger.info("search_complete session_id=%s query_chars=%s top_k=%s results=%s", session_id, len(query), top_k, len(results))
    return {"results": results}

@app.post("/quiz/")
async def quiz(request: Request):
    """Generate quiz from uploaded content."""
    data = await request.json()
    session_id = data.get("session_id")
    prompt = data.get("prompt") or "Quiz me on the uploaded material."
    context = retrieve_full_session_context(session_id)
    if not context:
        logger.warning("quiz_no_context session_id=%s", session_id)
        return {"quiz": [], "response": "Upload study files first so I can build a quiz from your material.", "sources": []}
    logger.info("quiz_begin session_id=%s prompt_chars=%s context_chunks=%s", session_id, len(prompt), len(context))
    try:
        chat_result = qwen_llm.chat_with_usage(prompt, context=context, study_task="quiz")
    except Exception:
        logger.exception("quiz_llm_failed session_id=%s context_chunks=%s", session_id, len(context))
        raise
    response = chat_result["response"]
    logger.info("quiz_complete session_id=%s response_chars=%s", session_id, len(response))
    return {"quiz": response, "response": response, "usage": chat_result["usage"], "sources": [public_source(source) for source in context]}

@app.post("/summarize/")
async def summarize(request: Request):
    """Summarize uploaded content or specific sections."""
    data = await request.json()
    session_id = data.get("session_id")
    prompt = data.get("prompt") or "Summarize the uploaded material."
    top_k = int(data.get("top_k", DEFAULT_RAG_TOP_K))
    context, _study_task = retrieve_study_context(session_id, prompt, top_k)
    if not context:
        logger.warning("summarize_no_context session_id=%s prompt_chars=%s", session_id, len(prompt))
        return {"summary": "Upload study files first so I can summarize your material.", "sources": []}
    logger.info("summarize_begin session_id=%s prompt_chars=%s context_chunks=%s", session_id, len(prompt), len(context))
    try:
        chat_result = qwen_llm.chat_with_usage(prompt, context=context, study_task="summary")
    except Exception:
        logger.exception("summarize_llm_failed session_id=%s context_chunks=%s", session_id, len(context))
        raise
    response = chat_result["response"]
    logger.info("summarize_complete session_id=%s response_chars=%s", session_id, len(response))
    return {"summary": response, "response": response, "usage": chat_result["usage"], "sources": [public_source(source) for source in context]}

@app.post("/video-link/")
async def video_link(request: Request):
    """Find and link to specific times in videos."""
    data = await request.json()
    session_id = data.get("session_id")
    prompt = data.get("prompt") or data.get("query") or ""
    context = retrieve_full_session_context(session_id)
    if not context:
        logger.warning("video_link_no_context session_id=%s prompt_chars=%s", session_id, len(prompt))
        return {"links": [], "response": "I could not find that topic in the uploaded material.", "sources": []}
    logger.info("video_link_begin session_id=%s prompt_chars=%s context_chunks=%s", session_id, len(prompt), len(context))
    try:
        chat_result = qwen_llm.chat_with_usage(prompt, context=context, study_task="find")
    except Exception:
        logger.exception("video_link_llm_failed session_id=%s context_chunks=%s", session_id, len(context))
        raise
    response = chat_result["response"]
    logger.info("video_link_complete session_id=%s response_chars=%s", session_id, len(response))
    return {"links": [public_source(source) for source in context], "response": response, "usage": chat_result["usage"], "sources": [public_source(source) for source in context]}
